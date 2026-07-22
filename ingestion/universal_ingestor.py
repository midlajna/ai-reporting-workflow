"""
universal_ingestor.py
Single entry point for ingesting ANY supported file type. Detects the
file extension and routes to the right parser, normalizing the result
into one of two buckets:

    {"type": "tabular", "source": <filename>, "data": <DataFrame>}
    {"type": "text",    "source": <filename>, "data": <str>}

This is what lets main.py accept a folder of mixed files instead of
fixed --csv / --doc arguments.
"""
import json
from pathlib import Path
import pandas as pd

from ingestion.tabular_ingestor import load_tabular_data
from ingestion.document_ingestor import load_document_text

TABULAR_EXTENSIONS = {".csv", ".xlsx", ".xls", ".json"}
TEXT_EXTENSIONS = {".pdf", ".txt", ".docx", ".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg"}


def load_json_as_table(filepath: str) -> pd.DataFrame:
    """Load a local JSON file (list of records, or {'data': [...]}) as a DataFrame."""
    with open(filepath, "r", encoding="utf-8") as f:
        payload = json.load(f)

    if isinstance(payload, dict):
        # try to find the first list value in the dict (common API shape)
        list_values = [v for v in payload.values() if isinstance(v, list)]
        payload = list_values[0] if list_values else [payload]

    return pd.json_normalize(payload)


def load_docx_text(filepath: str) -> str:
    """Extract all paragraph text from a Word document."""
    from docx import Document
    doc = Document(filepath)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    # also pull text out of any tables in the doc
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def load_pptx_text(filepath: str) -> str:
    """Extract all text from slides (titles, body text, notes)."""
    from pptx import Presentation
    prs = Presentation(filepath)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_text.append(shape.text_frame.text.strip())
        if slide_text:
            parts.append(f"[Slide {i}] " + " / ".join(slide_text))
    return "\n".join(parts)


def load_image_text_ocr(filepath: str) -> str:
    """Extract text from an image via OCR (tesseract)."""
    import pytesseract
    from PIL import Image
    img = Image.open(filepath)
    return pytesseract.image_to_string(img)


def ingest_any_file(filepath: str) -> dict:
    """Detect file type and return a normalized ingestion result."""
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"No such file: {filepath}")
    ext = path.suffix.lower()

    if ext in {".csv", ".xlsx", ".xls"}:
        return {"type": "tabular", "source": path.name, "data": load_tabular_data(filepath)}

    if ext == ".json":
        return {"type": "tabular", "source": path.name, "data": load_json_as_table(filepath)}

    if ext in {".pdf", ".txt"}:
        return {"type": "text", "source": path.name, "data": load_document_text(filepath)}

    if ext == ".docx":
        return {"type": "text", "source": path.name, "data": load_docx_text(filepath)}

    if ext == ".pptx":
        return {"type": "text", "source": path.name, "data": load_pptx_text(filepath)}

    if ext in IMAGE_EXTENSIONS:
        return {"type": "text", "source": path.name, "data": load_image_text_ocr(filepath)}

    raise ValueError(
        f"Unsupported file type: {ext}. "
        f"Supported: {sorted(TABULAR_EXTENSIONS | TEXT_EXTENSIONS | IMAGE_EXTENSIONS)}"
    )


def ingest_folder(folder_path: str) -> list:
    """Ingest every supported file in a folder. Skips unsupported files with a warning."""
    results = []
    supported = TABULAR_EXTENSIONS | TEXT_EXTENSIONS | IMAGE_EXTENSIONS
    for file in sorted(Path(folder_path).iterdir()):
        if file.is_dir():
            continue
        if file.suffix.lower() not in supported:
            print(f"  [skip] {file.name} — unsupported type")
            continue
        try:
            results.append(ingest_any_file(str(file)))
            print(f"  [ok]   {file.name} -> {results[-1]['type']}")
        except Exception as e:
            print(f"  [fail] {file.name} — {e}")
    return results


if __name__ == "__main__":
    results = ingest_folder("../sample_data")
    for r in results:
        print(r["source"], r["type"])
