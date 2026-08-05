"""
document_ingestor.py
Extracts plain text from unstructured documents.
Supports: .pdf, .txt, .docx, .pptx, and images (.png/.jpg/.jpeg) via OCR.
"""

from pathlib import Path
from typing import List
import re

SUPPORTED_DOCS = {".pdf", ".txt", ".docx", ".pptx", ".png", ".jpg", ".jpeg"}


def _extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _extract_pdf(path: Path) -> str:
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(path) as doc:
            for page in doc.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
        return "\n".join(texts)
    except Exception as e:
        print(f"[document_ingestor] PDF extract failed for {path.name}: {e}")
        return ""


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        print(f"[document_ingestor] DOCX extract failed for {path.name}: {e}")
        return ""


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        print(f"[document_ingestor] PPTX extract failed for {path.name}: {e}")
        return ""


def _extract_image_ocr(path: Path) -> str:
    try:
        from PIL import Image
        import pytesseract
        img = Image.open(path)
        return pytesseract.image_to_string(img)
    except Exception as e:
        print(f"[document_ingestor] OCR failed for {path.name}: {e}")
        return ""


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".txt":
        return _extract_txt(path)
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix in {".png", ".jpg", ".jpeg"}:
        return _extract_image_ocr(path)
    return ""


def load_documents(files: List[Path]) -> str:
    """
    Extract and concatenate text from all supported document/image files.
    Returns a single large string ready for summarization.
    """
    chunks = []
    for f in files:
        if f.suffix.lower() not in SUPPORTED_DOCS:
            continue
        text = extract_text(f)
        if text.strip():
            chunks.append(f"--- Source: {f.name} ---\n{text.strip()}")
    return "\n\n".join(chunks)
